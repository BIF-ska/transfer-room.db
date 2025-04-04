from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]  # 0 is the layout for the title slide
slide = prs.slides.add_slide(slide_layout)

# Set title and subtitle
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Automated PowerPoint Presentation"
subtitle.text = "Created using Python"

# Content slide
slide_layout = prs.slide_layouts[1]  # 1 is the layout for a title and content slide
slide = prs.slides.add_slide(slide_layout)

# Set title and content for the second slide
title = slide.shapes.title
title.text = "Introduction to Python"
content = slide.shapes.placeholders[1]
content.text = "fdgdgdfg""
# Add another slide with bullet points
slide_layout = prs.slide_layouts[1]  # You can choose the layout that suits you
slide = prs.slides.add_slide(slide_layout)

# Set title and bullet points
title = slide.shapes.title
title.text = "Key Features of Python"
content = slide.shapes.placeholders[1]
content.text = "• Easy syntax\n• Extensive libraries\n• Cross-platform\n• Open-source"

# Save the presentation
prs.save("automated_presentation.pptx")
